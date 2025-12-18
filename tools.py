import os
import json
import base64
import difflib
import re
import hashlib
import zipfile
import io
from typing import Any, Dict, List, Union
from datetime import datetime

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from utils import get_sandbox_path, is_vlm


# --- Tool Functions ---

def list_files(sandbox_id: str, max_depth: int = None) -> List[str]:
    """
    List all files in the sandbox directory.

    Args:
        sandbox_id: The sandbox ID
        max_depth: Maximum depth to traverse (None = unlimited, 1 = top-level only, etc.)
    """
    sandbox_path = get_sandbox_path(sandbox_id)
    output_files = []
    if not os.path.exists(sandbox_path):
         return ["Sandbox directory does not exist yet."]

    for root, dirs, files in os.walk(sandbox_path):
        # Calculate current depth relative to sandbox_path
        rel_root = os.path.relpath(root, sandbox_path)
        if rel_root == '.':
            current_depth = 0
        else:
            current_depth = rel_root.count(os.sep) + 1

        # Skip directories beyond max_depth
        if max_depth is not None and current_depth >= max_depth:
            dirs[:] = []  # Don't traverse deeper

        for file in files:
            file_path = os.path.join(root, file)
            rel_path = os.path.relpath(file_path, sandbox_path)
            if '.git' in rel_path or rel_path.startswith('.git/') or rel_path == 'conversation.json' or rel_path == '.context_cache':
                continue
            if os.path.isfile(file_path):
                output_files.append(rel_path)

    if len(output_files) == 0:
        return ["No files found in this sandbox."]
    return output_files


def read_file(sandbox_id: str, file_name: str, model_name: str = None, start_line: int = 1, end_line: int = -1) -> str:
    """
    Read a file from the sandbox directory, optionally specifying a line range.
    
    Args:
        sandbox_id: The sandbox ID
        file_name: The name of the file
        model_name: Optional model name for VLM logic
        start_line: 1-indexed, inclusive (default 1)
        end_line: 1-indexed, inclusive (default -1 for everything)
    """
    from file_preprocessor import get_converted_file_path
    import itertools

    sandbox_path = get_sandbox_path(sandbox_id)
    file_path = os.path.join(sandbox_path, file_name)
    if not os.path.exists(file_path):
        return f"Error: File {file_name} not found"

    # Check if this is a PDF that has been preprocessed
    if file_name.lower().endswith(".pdf"):
        # Determine expected format type
        format_type = "image" if is_vlm(model_name) else "text"
        converted_path = get_converted_file_path(file_path, format_type)

        # If converted file exists, read from it
        if os.path.exists(converted_path):
            try:
                with open(converted_path, 'r', encoding='utf-8') as f:
                    content = f.read()

                # If it's image format, the content is JSON
                if format_type == "image":
                    # Validate it's proper JSON
                    data = json.loads(content)
                    if data.get("__type__") == "image":
                        return content
                else:
                    # For text content, we can respect line numbers if it's just a raw text dump
                    # But if it's unstructured, slicing lines might be enough.
                    lines = content.splitlines(keepends=True)
                    
                    # Convert 1-indexed start/end to 0-indexed slices
                    start_idx = max(0, start_line - 1)
                    if end_line == -1:
                        target_lines = lines[start_idx:]
                    else:
                        target_lines = lines[start_idx:end_line]
                    
                    return "".join(target_lines)

            except Exception as e:
                # Fall through to read original if converted file is corrupted
                pass
        
        # If no converted file, return error suggesting preprocessing should have happened
        return f"Error: PDF file {file_name} has not been preprocessed. Please ensure file preprocessing runs before agent loop."

    # DOCX Support
    if file_name.lower().endswith(".docx"):
        if not docx:
            return "Error: python-docx library not available."
        try:
            doc = docx.Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
            lines = content.splitlines(keepends=True)
            start_idx = max(0, start_line - 1)
            stop_idx = None if end_line == -1 else end_line
            return "".join(lines[start_idx:stop_idx])
        except Exception as e:
            return f"Error reading DOCX: {e}"

    # PPTX Support
    if file_name.lower().endswith(".pptx"):
        if not Presentation:
            return "Error: python-pptx library not available."
        try:
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            content = "\n".join(text_runs)
            lines = content.splitlines(keepends=True)
            start_idx = max(0, start_line - 1)
            stop_idx = None if end_line == -1 else end_line
            return "".join(lines[start_idx:stop_idx])
        except Exception as e:
            return f"Error reading PPTX: {e}"

    try:
        # Optimized reading for large files
        if start_line == 1 and end_line == -1:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        else:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                # Use itertools.islice to efficiently skip to start_line
                # islice(iterable, start, stop) -> start is inclusive, stop is exclusive (0-indexed)
                # start_line is 1-indexed. so start=start_line-1
                # end_line is 1-indexed inclusive. so stop=end_line
                
                start_idx = max(0, start_line - 1)
                stop_idx = None if end_line == -1 else end_line
                
                # Careful with islice arguments: islice(f, start, stop)
                # If start_idx is 0, islice(f, stop_idx) returns first stop_idx lines
                # If start_idx > 0, islice(f, start_idx, stop_idx) 
                
                target_lines = itertools.islice(f, start_idx, stop_idx)
                return "".join(target_lines)

    except Exception as e:
         return f"Error reading file: {e}"


def write_to_file(sandbox_id: str, file_name: str, content: str) -> str:
    """Write content to a file in the sandbox."""
    sandbox_path = get_sandbox_path(sandbox_id)
    file_path = os.path.join(sandbox_path, file_name)
    try:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)

        # DOCX Support
        if file_name.lower().endswith(".docx"):
            if not docx:
                return "Error: python-docx library not available."
            doc = docx.Document()
            for line in content.splitlines():
                doc.add_paragraph(line)
            doc.save(file_path)
            return f"File {file_name} saved successfully as DOCX."

        # PPTX Support
        if file_name.lower().endswith(".pptx"):
            if not Presentation:
                return "Error: python-pptx library not available."
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Simple text box for content
            from pptx.util import Inches
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = content
            
            prs.save(file_path)
            return f"File {file_name} saved successfully as PPTX."

        with open(file_path, 'w') as f:
            f.write(content)
        return f"File {file_name} saved successfully."
    except Exception as e:
        return f"Error saving file: {e}"


def append_to_file(sandbox_id: str, file_name: str, content: str) -> str:
    """Append content to a file in the sandbox."""
    sandbox_path = get_sandbox_path(sandbox_id)
    file_path = os.path.join(sandbox_path, file_name)
    try:
        os.makedirs(os.path.dirname(file_path), exist_ok=True)

        # DOCX Support
        if file_name.lower().endswith(".docx"):
            if not docx:
                return "Error: python-docx library not available."
            if os.path.exists(file_path):
                doc = docx.Document(file_path)
            else:
                doc = docx.Document()
            for line in content.splitlines():
                doc.add_paragraph(line)
            doc.save(file_path)
            return f"Content appended to {file_name} (DOCX)."

        # PPTX Support
        if file_name.lower().endswith(".pptx"):
            if not Presentation:
                return "Error: python-pptx library not available."
            if os.path.exists(file_path):
                prs = Presentation(file_path)
            else:
                prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            from pptx.util import Inches
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = content
            
            prs.save(file_path)
            return f"Content appended to {file_name} (PPTX)."

        with open(file_path, 'a') as f:
            f.write(content)
        return f"Content appended to {file_name}."
    except Exception as e:
        return f"Error appending to file: {e}"


def delete_file(sandbox_id: str, file_name: str) -> str:
    """Delete a file in the sandbox."""
    sandbox_path = get_sandbox_path(sandbox_id)
    file_path = os.path.join(sandbox_path, file_name)
    try:
        if os.path.exists(file_path):
            os.remove(file_path)
            return f"File {file_name} deleted."
        return f"File {file_name} not found."
    except Exception as e:
        return f"Error deleting file: {e}"


def edit_file(sandbox_id: str, file_path: str, edits: List[Dict[str, str]], dry_run: bool = False, options: Dict[str, Any] = None) -> Union[Dict[str, Any], str]:
    """Make selective edits to files while preserving formatting."""
    
    def normalize_line_endings(text: str) -> str:
        return text.replace("\r\n", "\n")

    def normalize_whitespace(text: str) -> str:
        result = re.sub(r"[ \t]+", " ", text)
        result = "\n".join(line.strip() for line in result.split("\n"))
        return result

    def get_line_indentation(line: str) -> str:
        match = re.match(r"^(\s*)", line)
        return match.group(1) if match else ""

    def preserve_indentation(old_text: str, new_text: str) -> str:
        if ("- " in new_text or "* " in new_text) and ("- " in old_text or "* " in old_text):
            return new_text
        old_lines = old_text.split("\n")
        new_lines = new_text.split("\n")
        if not old_lines or not new_lines:
            return new_text
        base_indent = get_line_indentation(old_lines[0]) if old_lines and old_lines[0].strip() else ""
        old_indents = {i: get_line_indentation(line) for i, line in enumerate(old_lines) if line.strip()}
        new_indents = {i: get_line_indentation(line) for i, line in enumerate(new_lines) if line.strip()}
        first_new_indent_len = len(new_indents.get(0, "")) if new_indents else 0
        result_lines = []
        for i, new_line in enumerate(new_lines):
            if not new_line.strip():
                result_lines.append("")
                continue
            new_indent = new_indents.get(i, "")
            if i < len(old_lines) and i in old_indents:
                target_indent = old_indents[i]
            elif i == 0:
                target_indent = base_indent
            elif first_new_indent_len > 0:
                curr_indent_len = len(new_indent)
                indent_diff = max(0, curr_indent_len - first_new_indent_len)
                target_indent = base_indent
                for prev_i in range(i - 1, -1, -1):
                    if prev_i in old_indents and prev_i in new_indents:
                        prev_old = old_indents[prev_i]
                        prev_new = new_indents[prev_i]
                        if len(prev_new) <= curr_indent_len:
                            relative_spaces = curr_indent_len - len(prev_new)
                            target_indent = prev_old + " " * relative_spaces
                            break
            else:
                target_indent = new_indent
            result_lines.append(target_indent + new_line.lstrip())
        return "\n".join(result_lines)

    def create_unified_diff(original: str, modified: str, file_path: str) -> str:
        original_lines = original.splitlines(True)
        modified_lines = modified.splitlines(True)
        diff_lines = difflib.unified_diff(original_lines, modified_lines, fromfile=f"a/{file_path}", tofile=f"b/{file_path}", lineterm="")
        return "".join(diff_lines)

    def find_exact_match(content: str, pattern: str):
        if pattern in content:
            lines_before = content[: content.find(pattern)].count("\n")
            line_count = pattern.count("\n") + 1
            return True, lines_before, line_count
        return False, -1, 0

    full_sandbox_path = get_sandbox_path(sandbox_id)
    full_file_path = os.path.join(full_sandbox_path, file_path)

    if not file_path or not isinstance(file_path, str):
        return json.dumps({"success": False, "error": f"File path must be a non-empty string, got {type(file_path)}"})
    if not isinstance(edits, list) or not edits:
        return json.dumps({"success": False, "error": "Edits must be a non-empty list"})
    if not os.path.isfile(full_file_path):
        return json.dumps({"success": False, "error": f"File not found: {file_path}"})

    normalized_edits = []
    for i, edit in enumerate(edits):
        if not isinstance(edit, dict):
            return json.dumps({"success": False, "error": f"Edit #{i} must be a dictionary"})
        if "old_text" not in edit or "new_text" not in edit:
            return json.dumps({"success": False, "error": f"Edit #{i} missing old_text or new_text"})
        normalized_edits.append({"old_text": edit["old_text"], "new_text": edit["new_text"]})

    preserve_indent = options.get("preserve_indentation", True) if options else True
    normalize_ws = options.get("normalize_whitespace", True) if options else True

    try:
        if file_path.lower().endswith(".docx"):
            if not docx:
                return json.dumps({"success": False, "error": "python-docx library not available."})
            doc = docx.Document(full_file_path)
            original_content = "\n".join([para.text for para in doc.paragraphs])
        elif file_path.lower().endswith(".pptx"):
            if not Presentation:
                return json.dumps({"success": False, "error": "python-pptx library not available."})
            prs = Presentation(full_file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            original_content = "\n".join(text_runs)
        else:
            with open(full_file_path, "r", encoding="utf-8") as f:
                original_content = f.read()
    except Exception as e:
        return json.dumps({"success": False, "error": f"Error reading file: {str(e)}"})

    match_results = []
    changes_made = False
    modified_content = normalize_line_endings(original_content)
    
    for i, edit in enumerate(normalized_edits):
        old = normalize_line_endings(edit["old_text"])
        new = normalize_line_endings(edit["new_text"])
        
        if normalize_ws:
            # Simple exact match for now, as existing code does
            pass 

        if old in modified_content:
             modified_content = modified_content.replace(old, new, 1)
             changes_made = True
             match_results.append({"edit_index": i, "match_type": "exact"})
        else:
             match_results.append({"edit_index": i, "match_type": "failed"})

    if not changes_made:
         return json.dumps({"success": True, "message": "No changes made"})

    diff = create_unified_diff(original_content, modified_content, file_path)
    
    if not dry_run:
        try:
            if file_path.lower().endswith(".docx"):
                doc = docx.Document()
                for line in modified_content.splitlines():
                    doc.add_paragraph(line)
                doc.save(full_file_path)
            elif file_path.lower().endswith(".pptx"):
                prs = Presentation()
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                from pptx.util import Inches
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
                txBox.text_frame.text = modified_content
                prs.save(full_file_path)
            else:
                with open(full_file_path, "w", encoding="utf-8") as f:
                    f.write(modified_content)
        except Exception as e:
             return json.dumps({"success": False, "error": f"Error writing file: {str(e)}"})

    return json.dumps({"success": True, "diff": diff})


def get_folder_structure(sandbox_id: str, max_depth: int = 2) -> str:
    """Returns a tree-like string representation of the folder structure."""
    sandbox_path = get_sandbox_path(sandbox_id)
    if not os.path.exists(sandbox_path):
        return "Sandbox directory does not exist."
    
    output = []
    
    def add_to_tree(path, prefix="", depth=0):
        if depth > max_depth:
            return
        
        try:
            entries = sorted(os.listdir(path))
        except PermissionError:
            return

        entries = [e for e in entries if e != '.git' and e != '.context_cache' and e != 'conversation.json']
        
        for i, entry in enumerate(entries):
            is_last = (i == len(entries) - 1)
            full_path = os.path.join(path, entry)
            is_dir = os.path.isdir(full_path)
            
            connector = "└── " if is_last else "├── "
            output.append(f"{prefix}{connector}{entry}{'/' if is_dir else ''}")
            
            if is_dir:
                extension = "    " if is_last else "│   "
                add_to_tree(full_path, prefix + extension, depth + 1)
                
    output.append(".")
    add_to_tree(sandbox_path)
    return "\n".join(output)


def search_files(sandbox_id: str, pattern: str) -> List[str]:
    """Search for files in the sandbox matching a glob pattern."""
    import fnmatch
    sandbox_path = get_sandbox_path(sandbox_id)
    if not os.path.exists(sandbox_path):
        return ["Sandbox directory does not exist."]
    
    matches = []
    
    for root, dirs, files in os.walk(sandbox_path):
        # Skip git
        if '.git' in dirs:
            dirs.remove('.git')
            
        for name in files:
            if fnmatch.fnmatch(name, pattern):
                rel_path = os.path.relpath(os.path.join(root, name), sandbox_path)
                matches.append(rel_path)
    
    if not matches:
        return [f"No files found matching pattern '{pattern}'"]
    
    return matches


def search_content(sandbox_id: str, query: str, case_sensitive: bool = False) -> str:
    """Search for text content within files."""
    sandbox_path = get_sandbox_path(sandbox_id)
    if not os.path.exists(sandbox_path):
        return "Sandbox directory does not exist."
    
    results = []
    
    for root, dirs, files in os.walk(sandbox_path):
        if '.git' in dirs:
             dirs.remove('.git')
             
        for file in files:
            file_path = os.path.join(root, file)
            rel_path = os.path.relpath(file_path, sandbox_path)
            
            # Skip likely binary files or large files to prevent freezing
            skip_exts = ['.pyc', '.git', '.png', '.jpg', '.jpeg', '.pdf', '.zip', '.exe','conversation.json','.context_cache']
            if any(file.lower().endswith(ext) for ext in skip_exts):
                continue
                
            try:
                if file.lower().endswith(".docx"):
                    if not docx: continue
                    doc = docx.Document(file_path)
                    lines = [para.text for para in doc.paragraphs]
                elif file.lower().endswith(".pptx"):
                    if not Presentation: continue
                    prs = Presentation(file_path)
                    lines = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                lines.append(shape.text)
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        lines = f.readlines()
                    
                for i, line in enumerate(lines):
                    if case_sensitive:
                        match = query in line
                    else:
                        match = query.lower() in line.lower()
                        
                    if match:
                        # Truncate line if too long
                        content = line.strip()
                        if len(content) > 200:
                            content = content[:200] + "..."
                        results.append(f"{rel_path}:{i+1}: {content}")
                        if len(results) > 100:  # Global limit
                             break
            except Exception:
                pass  # Skip files we can't read
            
            if len(results) > 100:
                results.append("... (results truncated)")
                break
                
    if not results:
        return f"No matches found for '{query}'"
        
    return "\n".join(results)


def get_file_info(sandbox_id: str, file_path: str) -> str:
    """Get metadata about a specific file."""
    sandbox_path = get_sandbox_path(sandbox_id)
    full_path = os.path.join(sandbox_path, file_path)
    
    if not os.path.exists(full_path):
        return f"File not found: {file_path}"
        
    try:
        stat = os.stat(full_path)
        is_dir = os.path.isdir(full_path)
        
        info = {
            "name": os.path.basename(file_path),
            "type": "directory" if is_dir else "file",
            "size_bytes": stat.st_size,
            "created": datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S'),
            "modified": datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
        }
        
        if not is_dir:
            # Count lines for text files
            try:
                with open(full_path, 'r', encoding='utf-8') as f:
                    info["lines"] = sum(1 for _ in f)
            except:
                info["lines"] = "N/A (binary or non-utf8)"
                
        return json.dumps(info, indent=2)
    except Exception as e:
        return f"Error getting file info: {e}"


        return f"Error getting file info: {e}"


def import_outlook_emails(sandbox_id: str, query: str = None, received_after: str = None) -> str:
    """
    Import emails from Outlook based on a grep-like query and optional time boundary.
    Only works on Windows.
    """
    try:
        import win32com.client
    except ImportError:
        return "Error: Outlook import tool is only available on Windows systems with pywin32 installed."

    sandbox_path = get_sandbox_path(sandbox_id)
    if not os.path.exists(sandbox_path):
        return "Error: Sandbox directory does not exist."

    processed_count = 0
    saved_paths = []
    
    # Parse date if provided
    filter_date = None
    outlook_date_filter = ""
    if received_after:
        try:
            # Parse YYYY-MM-DD
            dt = datetime.strptime(received_after, "%Y-%m-%d")
            filter_date = dt
            # Outlook Restrict filter format: [ReceivedTime] >= 'MM/DD/YYYY 00:00 AM'
            outlook_date_filter = f"[ReceivedTime] >= '{dt.strftime('%m/%d/%Y')} 00:00 AM'"
        except ValueError:
            return "Error: received_after must be in YYYY-MM-DD format."
            
    from file_preprocessor import convert_pdf_to_text

    def sanitize_filename(name: str) -> str:
        """Sanitize string for usage as filename."""
        return re.sub(r'[<>:"/\\|?*]', '_', str(name)).strip()[:50]

    def _save_outlook_message(message, folder_name_source: str) -> str:
        """Helper to save an Outlook message and its attachments."""
        nonlocal processed_count
        try:
            # Filter for emails specifically if possible, but handle others
            # olMail=43, olAppointment=26, olMeetingRequest=53
            item_class = getattr(message, 'Class', 43)
            
            subj = getattr(message, 'Subject', 'No Subject') or 'No Subject'
            body = getattr(message, 'Body', '') or ''
            html_body = getattr(message, 'HTMLBody', '') or ''
            
            # Use EntryID if available for unique naming, fallback to hash
            entry_id = getattr(message, 'EntryID', None)
            if entry_id:
                unique_id = hashlib.sha256(entry_id.encode('utf-8')).hexdigest()
            else:
                unique_id = hashlib.sha256((subj + body[:100]).encode('utf-8', errors='ignore')).hexdigest()
            
            sender_name = "Unknown"
            try:
                if item_class == 43: # MailItem
                    sender_name = getattr(message, 'SenderName', '') or getattr(message, 'SenderEmailAddress', 'Unknown')
                elif item_class == 26: # AppointmentItem
                    sender_name = getattr(message, 'Organizer', 'Unknown')
            except: pass

            folder_name = f"{sanitize_filename(sender_name)}_{sanitize_filename(subj)}_{unique_id[:8]}"
            save_folder = os.path.join(sandbox_path, "mail", folder_name)
            
            if os.path.exists(save_folder):
                return None

            os.makedirs(save_folder, exist_ok=True)

            meta = {
                "id": unique_id,
                "EntryID": entry_id,
                "Subject": subj,
                "Body": body,
                "HTMLBody": html_body[:5000] if html_body else "", # Truncate HTML body for meta
                "ReceivedTime": str(getattr(message, 'ReceivedTime', '')),
                "Sender": sender_name,
                "To": getattr(message, 'To', ''),
                "FolderName": folder_name_source,
                "ItemClass": item_class
            }
            
            with open(os.path.join(save_folder, "email_data.json"), 'w', encoding='utf-8') as f:
                json.dump(meta, f, indent=2)

            if hasattr(message, 'Attachments'):
                for attachment in message.Attachments:
                    try:
                        file_path = os.path.join(save_folder, attachment.FileName)
                        attachment.SaveAsFile(file_path)
                        
                        if attachment.FileName.lower().endswith('.pdf'):
                            try:
                                content, format_type = convert_pdf_to_text(file_path)
                                ext = ".txt" if format_type == "text" else ".json"
                                with open(file_path + ext, 'w', encoding='utf-8') as f:
                                    f.write(content)
                            except Exception: pass

                        if attachment.FileName.lower().endswith('.zip'):
                            try:
                                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                                    zip_ref.extractall(save_folder)
                            except Exception: pass 
                    except Exception: pass

            processed_count += 1
            return os.path.relpath(save_folder, sandbox_path)
        except Exception:
            return None

    try:
        outlook_app = win32com.client.Dispatch('Outlook.Application')
        namespace = outlook_app.GetNamespace('MAPI')
        
        # 1. Strategy: Windows Search Index (ADODB) - Fast!
        if query:
            try:
                conn = win32com.client.Dispatch("ADODB.Connection")
                conn.Open("Provider=Search.CollatorDSO;Extended Properties='Application=Windows';")
                
                # We fetch System.ItemUrl which contains the outlook:0000... protocol with EntryID
                sql = f"""
                    SELECT "System.ItemUrl" FROM SystemIndex 
                    WHERE System.Kind = 'email' 
                    AND CONTAINS(*, '"{query}"')
                """
                if received_after:
                    # Windows Search date format is YYYY/MM/DD
                    sql += f" AND System.DateModified >= '{received_after.replace('-', '/')}'"
                sql += " ORDER BY System.DateModified DESC"
                
                rs = conn.Execute(sql)[0]
                while not rs.EOF:
                    if processed_count >= 50: break
                    item_url = rs.Fields.Item(0).Value # e.g. outlook:00000000...
                    try:
                        # Extract entry ID from URL if it starts with outlook:
                        if item_url.startswith("outlook:"):
                            entry_id = item_url.split("outlook:")[1]
                            message = namespace.GetItemFromID(entry_id)
                            path = _save_outlook_message(message, "WindowsIndex")
                            if path: saved_paths.append(path)
                    except: pass
                    rs.MoveNext()
                rs.Close()
                conn.Close()
            except Exception:
                pass

        # 2. Strategy: DASL Filter (Outlook Native Search) - Medium Speed
        if processed_count < 10: # Only if ADODB didn't yield much
            targets = [namespace.GetDefaultFolder(6), namespace.GetDefaultFolder(5)]
            
            dasl_query = ""
            if query:
                dasl_query = (
                    f"@SQL=\"urn:schemas:httpmail:subject\" LIKE '%{query}%' OR "
                    f"\"urn:schemas:httpmail:textdescription\" LIKE '%{query}%' OR "
                    f"\"urn:schemas:httpmail:fromname\" LIKE '%{query}%'"
                )
            
            if received_after:
                try:
                    dt_obj = datetime.strptime(received_after, "%Y-%m-%d")
                    date_part = f"\"urn:schemas:httpmail:datereceived\" >= '{dt_obj.strftime('%m/%d/%Y')} 00:00 AM'"
                    if dasl_query:
                        dasl_query = f"@SQL=({dasl_query.replace('@SQL=', '')}) AND {date_part}"
                    else:
                        dasl_query = f"@SQL={date_part}"
                except: pass

            for folder in targets:
                if processed_count >= 50: break
                try:
                    items = folder.Items
                    if dasl_query:
                        items = items.Restrict(dasl_query)
                    items.Sort("[ReceivedTime]", True)
                    
                    for message in items:
                        if processed_count >= 50: break
                        path = _save_outlook_message(message, folder.Name)
                        if path: saved_paths.append(path)
                except Exception: continue

        # 3. Strategy: Full Recursive Fallback (Current) - Slowest
        if processed_count == 0:
            def process_folder_recursive(folder):
                for sub in folder.Folders:
                    if processed_count >= 50: return
                    process_folder_recursive(sub)
                
                items = folder.Items
                # Use crude restriction for dates if possible
                if received_after:
                    try:
                        dt_obj = datetime.strptime(received_after, "%Y-%m-%d")
                        local_filter = f"[ReceivedTime] >= '{dt_obj.strftime('%m/%d/%Y')} 00:00 AM'"
                        items = items.Restrict(local_filter)
                    except: pass
                
                for msg in items:
                    if processed_count >= 50: return
                    if query:
                        subj = getattr(msg, 'Subject', '') or ''
                        body = getattr(msg, 'Body', '') or ''
                        if not re.search(query, f"{subj} {body}", re.I): continue
                    
                    path = _save_outlook_message(msg, folder.Name)
                    if path: saved_paths.append(path)

            for account in namespace.Folders:
                if processed_count >= 50: break
                process_folder_recursive(account)

    except Exception as e:
        return f"Error connecting to Outlook: {e}"

    if processed_count == 0:
        return "No new emails found matching criteria."
    
    return f"Imported {processed_count} emails. Saved to: {', '.join(saved_paths[:5])}..."


# --- Tool Registry ---

TOOL_DEFINITIONS = [
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "List all files in the current working directory.",
            "parameters": {
                "type": "object",
                "properties": {
                },
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_folder_structure",
            "description": "Get a tree-like representation of the folders.",
            "parameters": {
                "type": "object",
                "properties": {
                    "max_depth": {"type": "integer", "description": "Max depth to traverse (default 2)."}
                },
                "required": []
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_files",
            "description": "Search for files matching a pattern (e.g. *.py).",
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {"type": "string", "description": "Glob pattern to search for."}
                },
                "required": ["pattern"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "search_content",
            "description": "Search for text content within files (grep-like).",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Text to search for."},
                    "case_sensitive": {"type": "boolean", "description": "Case sensitive match."}
                },
                "required": ["query"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "get_file_info",
            "description": "Get metadata (size, lines, time) for a file.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Relative path to the file."}
                },
                "required": ["file_path"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read a file from the current working directory. Use start_line/end_line for pagination.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_name": {"type": "string", "description": "The name of the file to read."},
                    "start_line": {"type": "integer", "description": "Line to start reading from (1-indexed).", "default": 1},
                    "end_line": {"type": "integer", "description": "Line to stop reading at (1-indexed, inclusive). -1 for end of file.", "default": -1}
                },
                "required": ["file_name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "write_to_file",
            "description": "Write content to a file (overwrites).",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_name": {"type": "string", "description": "The name of the file to write to."},
                    "content": {"type": "string", "description": "Content to write."}
                },
                "required": ["file_name", "content"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "append_to_file",
            "description": "Append content to a file.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_name": {"type": "string", "description": "The name of the file to append to."},
                    "content": {"type": "string", "description": "Content to append."}
                },
                "required": ["file_name", "content"]
            }
        }
    },
     {
        "type": "function",
        "function": {
            "name": "delete_file",
            "description": "Delete a file.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_name": {"type": "string", "description": "The name of the file to delete."}
                },
                "required": ["file_name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "edit_file",
            "description": "Edit a file using search and replace blocks.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the file."},
                    "edits": {
                        "type": "array", 
                        "items": {
                            "type": "object",
                            "properties": {
                                "old_text": {"type": "string"},
                                "new_text": {"type": "string"}
                            },
                            "required": ["old_text", "new_text"]
                        }
                    }
                },
                "required": ["file_path", "edits"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "import_outlook_emails",
            "description": "Import emails from Outlook (Windows only) using a grep-like query and time boundary.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Regex pattern to search in Subject, Body, Sender, and To fields."
                    },
                    "received_after": {
                        "type": "string",
                        "description": "Optional: Filter emails received after this date (YYYY-MM-DD)."
                    }
                },
                "required": ["query"]
            }
        }
    }
]


# --- Tool Execution ---

def execute_tool(name: str, args: Dict[str, Any]) -> str:
    """Execute a tool by name with arguments."""
    try:
        # sandbox_id is injected into args by agent.py
        if name == "list_files":
            return str(list_files(args.get("sandbox_id"), args.get("max_depth")))
        elif name == "get_folder_structure":
            return str(get_folder_structure(args.get("sandbox_id"), args.get("max_depth", 2)))
        elif name == "search_files":
            return str(search_files(args.get("sandbox_id"), args.get("pattern")))
        elif name == "search_content":
             return str(search_content(args.get("sandbox_id"), args.get("query"), args.get("case_sensitive", False)))
        elif name == "get_file_info":
             return str(get_file_info(args.get("sandbox_id"), args.get("file_path")))
        elif name == "read_file":
            return str(read_file(
                args.get("sandbox_id"), 
                args.get("file_name"), 
                args.get("model_name"),
                args.get("start_line", 1),
                args.get("end_line", -1)
            ))
        elif name == "write_to_file":
            return str(write_to_file(args.get("sandbox_id"), args.get("file_name"), args.get("content")))
        elif name == "append_to_file":
            return str(append_to_file(args.get("sandbox_id"), args.get("file_name"), args.get("content")))
        elif name == "delete_file":
            return str(delete_file(args.get("sandbox_id"), args.get("file_name")))
        elif name == "edit_file":
            return str(edit_file(
                args.get("sandbox_id"), 
                args.get("file_path"), 
                args.get("edits"), 
                args.get("dry_run", False), 
                args.get("options")
            ))
        elif name == "import_outlook_emails":
            return str(import_outlook_emails(args.get("sandbox_id"), args.get("query"), args.get("received_after")))
        else:
            return f"Error: Tool {name} not found."
    except Exception as e:
        return f"Error executing tool {name}: {e}"
