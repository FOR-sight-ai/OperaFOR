import os
import json
import zipfile
import shutil
import tempfile
import re
from pathlib import Path
from typing import Dict, List, Any, Optional

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
except ImportError:
    Presentation = None

class PPTXHandler:
    def __init__(self, file_path: str):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

    def get_inventory(self) -> Dict[str, Any]:
        """
        Extracts a structured inventory of text elements from the PPTX.
        Returns a dict: { slide_id: { shape_id: { properties, paragraphs: [...] } } }
        """
        if Presentation is None:
            return {"error": "python-pptx library not available"}

        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            return {"error": f"Could not open PPTX: {e}"}

        inventory = {}

        for slide_idx, slide in enumerate(prs.slides):
            slide_key = f"slide-{slide_idx + 1}"
            slide_inventory = {}
            
            # Helper to collect text from shapes (including groups)
            def collect_text_shapes(shapes, parent_inventory):
                for shape_idx, shape in enumerate(shapes):
                    if hasattr(shape, "shapes"): # GroupShape
                        collect_text_shapes(shape.shapes, parent_inventory)
                        continue
                        
                    if not hasattr(shape, "text_frame") or not shape.text_frame:
                        continue
                    
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue

                    shape_id = f"shape-{len(parent_inventory) + 1}"
                    
                    paragraphs = []
                    for p in shape.text_frame.paragraphs:
                        if not p.text.strip():
                            continue
                        para_data = {
                            "text": p.text,
                            "level": getattr(p, "level", 0),
                        }
                        if p.alignment:
                            para_data["alignment"] = str(p.alignment)
                        
                        # Add run info for even more detail
                        runs = []
                        for r in p.runs:
                            run_data = {"text": r.text}
                            if r.font.name: run_data["font"] = r.font.name
                            if r.font.size: run_data["size"] = r.font.size.pt
                            if r.font.bold: run_data["bold"] = True
                            runs.append(run_data)
                        
                        para_data["runs"] = runs
                        paragraphs.append(para_data)

                    shape_data = {
                        "name": shape.name,
                        "paragraphs": paragraphs
                    }
                    parent_inventory[shape_id] = shape_data

            collect_text_shapes(slide.shapes, slide_inventory)
            
            if slide_inventory:
                inventory[slide_key] = slide_inventory

        return inventory

    def replace_text(self, replacements: List[Dict[str, str]]) -> Dict[str, Any]:
        """
        Replaces text in-place using python-pptx.
        'replacements' is a list of {"old_text": "...", "new_text": "..."}
        """
        if Presentation is None:
            return {"success": False, "error": "python-pptx library not available"}

        if not replacements:
            return {"success": True, "message": "No replacements provided"}

        try:
            prs = Presentation(self.file_path)
            total_matches = 0

            def process_shapes(shapes):
                nonlocal total_matches
                for shape in shapes:
                    if hasattr(shape, "shapes"):
                        process_shapes(shape.shapes)
                        continue
                    
                    if not hasattr(shape, "text_frame") or not shape.text_frame:
                        continue
                    
                    for paragraph in shape.text_frame.paragraphs:
                        # Try to replace in the whole paragraph text if splits are an issue
                        # But setting paragraph.text clears all runs.
                        # If we want to preserve run formatting, we must do it run by run.
                        for run in paragraph.runs:
                            for rep in replacements:
                                old = rep.get("old_text")
                                new = rep.get("new_text")
                                if old and old in run.text:
                                    run.text = run.text.replace(old, new)
                                    total_matches += 1
                        
                        # Fallback for split runs: if the joined text matches but runs don't
                        # this is harder. For now, we stick to run-by-run to keep formatting.

            for slide in prs.slides:
                process_shapes(slide.shapes)

            if total_matches > 0:
                prs.save(self.file_path)
                return {"success": True, "message": f"Applied {total_matches} replacements."}
            else:
                return {"success": True, "message": "No matches found."}

        except Exception as e:
            return {"success": False, "error": str(e)}

    def _create_zip(self, source_dir: Path, output_file: Path):
        with zipfile.ZipFile(output_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(source_dir):
                for file in files:
                    file_path = Path(root) / file
                    arcname = file_path.relative_to(source_dir)
                    zipf.write(file_path, arcname)

def get_pptx_inventory(file_path: str) -> str:
    try:
        handler = PPTXHandler(file_path)
        inventory = handler.get_inventory()
        return json.dumps(inventory, indent=2)
    except Exception as e:
        return json.dumps({"error": str(e)})

def edit_pptx_inplace(file_path: str, replacements: List[Dict[str, str]]) -> str:
    try:
        handler = PPTXHandler(file_path)
        result = handler.replace_text(replacements)
        return json.dumps(result)
    except Exception as e:
        return json.dumps({"success": False, "error": str(e)})
